"""
Minimal touch-vs-click.pptx (English) + diagram images (Pillow).
Requires: pip install python-pptx pillow
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

ASSETS_DIR = Path(__file__).resolve().parent / "assets_touch_click"


def load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path(r"C:\Windows\Fonts\segoeui.ttf"),
        Path(r"C:\Windows\Fonts\arial.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
    ]
    for p in candidates:
        if p.is_file():
            try:
                return ImageFont.truetype(str(p), size=size)
            except OSError:
                continue
    return ImageFont.load_default()


def draw_arrow(draw: ImageDraw.ImageDraw, x1: int, y1: int, x2: int, y2: int, fill: str, width: int = 3) -> None:
    draw.line((x1, y1, x2, y2), fill=fill, width=width)
    # small arrowhead
    if x2 > x1:
        draw.polygon([(x2, y2), (x2 - 12, y2 - 6), (x2 - 12, y2 + 6)], fill=fill)


def build_figure_touch_timeline(path: Path) -> None:
    w, h = 2000, 940
    img = Image.new("RGB", (w, h), "#f8fafc")
    draw = ImageDraw.Draw(img)
    font_title = load_font(42)
    font_lg = load_font(26)
    font_sm = load_font(20)
    font_mono = load_font(22)
    font_bullet = load_font(19)

    draw.text((60, 36), "Touch events = full finger story (many steps over time)", fill="#0f172a", font=font_title)

    y = 220
    # Each tuple: x0, y0, x1, y1 (x1 > x0)
    boxes = [
        (80, y, 360, y + 88, "#7c3aed", "touchstart", "Finger down"),
        (400, y + 14, 560, y + 74, "#a78bfa", "touchmove", "…"),
        (590, y + 14, 750, y + 74, "#a78bfa", "touchmove", "…"),
        (780, y + 14, 980, y + 74, "#a78bfa", "touchmove", "slides"),
        (1010, y, 1280, y + 88, "#7c3aed", "touchend", "Finger up"),
        (1310, y + 10, 1880, y + 78, "#94a3b8", "touchcancel", "Optional if browser stops (e.g. scroll)"),
    ]
    for i, (x0, y0, x1, y1, fill, name, sub) in enumerate(boxes):
        draw.rounded_rectangle((x0, y0, x1, y1), radius=14, fill=fill, outline="#0f172a", width=2)
        tw = draw.textlength(name, font=font_mono)
        draw.text((x0 + (x1 - x0 - tw) / 2, y0 + 10), name, fill="white" if fill != "#94a3b8" else "#1e293b", font=font_mono)
        stw = draw.textlength(sub, font=font_sm)
        draw.text((x0 + (x1 - x0 - stw) / 2, y0 + 44), sub, fill="white" if fill != "#94a3b8" else "#334155", font=font_sm)
        if i < len(boxes) - 1:
            nx0 = boxes[i + 1][0]
            my = (y0 + y1) // 2
            draw_arrow(draw, x1 + 6, my, nx0 - 6, my, "#334155", 4)

    draw.text((60, 350), "→ time", fill="#64748b", font=font_lg)
    draw.line((60, 388, w - 60, 388), fill="#cbd5e1", width=4)

    draw.rounded_rectangle((60, 415, w - 60, h - 20), radius=16, fill="#fff", outline="#e2e8f0", width=2)
    draw.text((90, 432), "Use touch — concrete examples", fill="#0f172a", font=font_lg)
    draw.text(
        (90, 472),
        "You need continuous (x, y) while the finger moves — often use Pointer Events today (same idea + mouse).",
        fill="#334155",
        font=font_sm,
    )
    examples = [
        "• Slider / range thumb — track touchmove or pointermove for position",
        "• Map pan — finger down + many moves until up",
        "• Sortable list — drag handle (not the same as click-to-select row)",
        "• Signature pad / drawing canvas — many touchmove points per stroke",
        "• Pull-to-refresh — measure distance from touchstart to touchend / move",
        "• Game virtual joystick / on-screen D-pad",
        "• Image crop / resize handles at corners and edges",
    ]
    y_after = draw_lines(draw, 90, 505, examples, font_bullet, "#334155", gap=3)
    draw.text(
        (90, y_after + 10),
        "touchmove can fire tens or hundreds of times per gesture — that is normal, not a bug.",
        fill="#64748b",
        font=font_sm,
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, "PNG", optimize=True)


def draw_lines(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    lines: list[str],
    font: ImageFont.FreeTypeFont | ImageFont.ImageFont,
    fill: str,
    gap: int = 6,
) -> int:
    """Draw stacked lines; returns next y below last line."""
    for line in lines:
        draw.text((x, y), line, fill=fill, font=font)
        bbox = draw.textbbox((x, y), line, font=font)
        y = bbox[3] + gap
    return y


def build_figure_click(path: Path) -> None:
    w, h = 2000, 920
    img = Image.new("RGB", (w, h), "#f8fafc")
    draw = ImageDraw.Draw(img)
    font_title = load_font(42)
    font_lg = load_font(26)
    font_sm = load_font(20)
    font_bullet = load_font(19)

    draw.text((60, 36), "Click = one activation event (short signal)", fill="#0f172a", font=font_title)

    # Left inputs
    draw.rounded_rectangle((80, 200, 420, 320), radius=16, fill="#e0f2fe", outline="#0284c7", width=2)
    draw.text((110, 220), "Desktop", fill="#0369a1", font=font_lg)
    draw.text((110, 260), "mouse down + up\non same target", fill="#0c4a6e", font=font_sm)

    draw.rounded_rectangle((80, 360, 420, 500), radius=16, fill="#d1fae5", outline="#059669", width=2)
    draw.text((110, 380), "Phone / tablet", fill="#047857", font=font_lg)
    draw.text((110, 420), "tap (after touch\nsequence ends)", fill="#064e3b", font=font_sm)

    draw.text((480, 300), "Browser decides\n“this counts as\nactivation”", fill="#334155", font=font_lg)

    draw_arrow(draw, 430, 260, 700, 320, "#64748b", 4)
    draw_arrow(draw, 430, 430, 700, 360, "#64748b", 4)

    draw.rounded_rectangle((760, 240, 1180, 440), radius=20, fill="#059669", outline="#064e3b", width=3)
    tw = draw.textlength("click", font=load_font(56))
    draw.text((760 + (1180 - 760 - tw) / 2, 300), "click", fill="white", font=load_font(56))

    draw.text((1240, 280), "One DOM name:\nElement.addEventListener(\"click\", …)", fill="#0f172a", font=font_lg)

    draw.rounded_rectangle((60, 500, w - 60, h - 20), radius=16, fill="#fff", outline="#e2e8f0", width=2)
    draw.text((90, 520), "Use click — concrete examples", fill="#0f172a", font=font_lg)
    examples = [
        "• Save / Submit / Delete confirm on <button>",
        "• Follow a link (<a href>) or router navigation",
        "• Open / close modal — OK, Cancel, dismiss overlay",
        "• Hamburger / menu toggle, icon buttons (one-shot actions)",
        "• Accordion “click header to expand” (if not dragging)",
        "• Row select in a table when the row is not a drag handle",
        "• Anything that must work with keyboard (Space / Enter) on native <button>",
    ]
    draw_lines(draw, 90, 560, examples, font_bullet, "#334155", gap=4)

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, "PNG", optimize=True)


def build_figure_compare(path: Path) -> None:
    w, h = 2000, 820
    img = Image.new("RGB", (w, h), "#f8fafc")
    draw = ImageDraw.Draw(img)
    font_title = load_font(40)
    font_lg = load_font(24)
    font_sm = load_font(19)

    draw.text((60, 28), "Same tap on a phone — two layers of events", fill="#0f172a", font=font_title)

    draw.rounded_rectangle((60, 100, 200, 640), radius=12, fill="#e2e8f0", outline="#94a3b8", width=2)
    draw.text((95, 320), "one\ntap", fill="#475569", font=font_lg)

    draw_arrow(draw, 210, 360, 250, 360, "#64748b", 4)

    # Touch layer
    draw.text((260, 110), "Layer 1 — touch timeline (long)", fill="#5b21b6", font=font_lg)
    draw.rounded_rectangle((260, 150, 1920, 260), radius=18, fill="#ede9fe", outline="#7c3aed", width=3)
    fm = load_font(22)
    x = 320
    for label in ["touchstart", "touchmove", "touchmove", "…", "touchend"]:
        draw.text((x, 188), label, fill="#4c1d95", font=fm)
        x += 280

    # Click layer
    draw.text((260, 300), "Layer 2 — click (short, often after touch)", fill="#047857", font=font_lg)
    draw.rounded_rectangle((260, 340, 1920, 430), radius=18, fill="#d1fae5", outline="#059669", width=3)
    draw.text((1200, 365), "click", fill="#064e3b", font=load_font(36))

    draw.rounded_rectangle((60, 480, w - 60, 780), radius=16, fill="#fff7ed", outline="#fb923c", width=2)
    draw.text((90, 510), "Warning — double submit bug", fill="#9a3412", font=load_font(28))
    draw.text(
        (90, 560),
        "If your code runs the same action on touchend and on click, one tap can trigger twice.",
        fill="#7c2d12",
        font=font_lg,
    )
    draw.text(
        (90, 610),
        "Fix: handle the business logic in one place, or ignore the second event with a guard.",
        fill="#9a3412",
        font=font_sm,
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    img.save(path, "PNG", optimize=True)


def build_all_figures() -> dict[str, Path]:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    paths = {
        "touch": ASSETS_DIR / "diagram_touch_timeline.png",
        "click": ASSETS_DIR / "diagram_click.png",
        "compare": ASSETS_DIR / "diagram_compare_layers.png",
    }
    build_figure_touch_timeline(paths["touch"])
    build_figure_click(paths["click"])
    build_figure_compare(paths["compare"])
    return paths


def set_body(tf, size: int = 22):
    for p in tf.paragraphs:
        p.font.name = "Calibri"
        p.font.size = Pt(size)


def add_title(prs, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets(prs, title: str, lines: list[str], *, size: int = 22):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
    set_body(body, size)


def add_diagram_slide(prs: Presentation, title: str, image_path: Path, caption: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title
    left = Inches(0.65)
    top = Inches(1.05)
    width = Inches(12.0)
    slide.shapes.add_picture(str(image_path.resolve()), left, top, width=width)
    if caption:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(6.72), Inches(12.1), Inches(0.55))
        tf = box.text_frame
        tf.text = caption
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.name = "Calibri"
        tf.paragraphs[0].font.italic = True


def main():
    figures = build_all_figures()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title(
        prs,
        "Touch vs Click",
        "Short deck • Diagrams + live demo: demo/index.html",
    )

    add_bullets(
        prs,
        "Click",
        [
            "One DOM event: click.",
            "Fires when the browser treats a press-and-release as “activate this” (mouse, and often after a tap on phones).",
            "Good when you only need that single moment — like pressing a doorbell once.",
        ],
    )
    add_diagram_slide(
        prs,
        "Figure — what “click” means",
        figures["click"],
        "One event name in your code; browser merges mouse or tap into this signal when appropriate.",
    )

    add_bullets(
        prs,
        "Touch",
        [
            "Several DOM events: touchstart → touchmove (many times) → touchend or touchcancel.",
            "Describes the whole finger movement until it ends — a timeline, not one beep.",
            "Good when you care where the finger is while it moves (drag, draw, swipe).",
        ],
    )
    add_diagram_slide(
        prs,
        "Figure — touch event timeline",
        figures["touch"],
        "Gray box = optional; touchcancel happens when the browser gives up on tracking (e.g. scroll).",
    )

    add_bullets(
        prs,
        "Side by side",
        [
            "Touch = long story of the finger. Click = short “that counted as a button press” at the end (often).",
            "One tap on a phone can produce touch events first, then a click — same gesture, different layers.",
            "Do not run the same important logic twice on both touchend and click unless you guard it.",
        ],
        size=21,
    )
    add_diagram_slide(
        prs,
        "Figure — one tap, two layers",
        figures["compare"],
        "Same physical tap: first the touch timeline, then sometimes a click — not two separate user actions.",
    )

    add_bullets(
        prs,
        "Real apps — use click",
        [
            "Submit, OK, open dialog, follow link, toggle menu on a <button>.",
            "Anything that should also work with keyboard (Space/Enter) on native buttons.",
        ],
        size=24,
    )

    add_bullets(
        prs,
        "Real apps — use touch (or Pointer Events)",
        [
            "Slider thumb, map drag, sortable row drag, signature pad, game control.",
            "Custom gestures (swipe distance, long-press) where you read many move points.",
            "New code: Pointer Events often replace separate mouse + touch handlers for drags.",
        ],
        size=22,
    )

    add_bullets(
        prs,
        "Demo",
        [
            "Open demo/index.html (+ demo.js, demo.css). Tap button logs the full pointer+touch+mouse+click sequence from one press.",
            "Left: Tap button + one drag rail. Right: log.",
        ],
        size=22,
    )

    out = Path(__file__).resolve().parent / "touch-vs-click.pptx"
    prs.save(str(out))
    print(f"Wrote {out}")
    print(f"Figures: {ASSETS_DIR}")


if __name__ == "__main__":
    main()
